import io
import re

from lxml import etree
from pptx import Presentation

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
    "dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def _collapse(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def _paragraph_texts_from_shape(shape) -> list[str]:
    """One entry per paragraph, matching python-docx/pptx's own paragraph
    boundaries — NOT joined across paragraphs. This mirrors the fix in
    apps/contentcheck/js/parser.js where concatenating every run in a shape
    into one string glued unrelated bullets together (e.g. "VF"+"XYZ"+"XYZ"
    -> "VFXYZXYZ"), making them impossible to match against the other doc.
    """
    texts = []
    for paragraph in shape.text_frame.paragraphs:
        text = _collapse(paragraph.text)
        if text:
            texts.append(text)
    return texts


def _extract_diagram_texts(shape) -> list[str]:
    """SmartArt diagrams aren't exposed by python-pptx at all — the actual
    text lives in a separate ppt/diagrams/dataN.xml part referenced by a
    relationship id on the graphicFrame. Without this, content redesigned
    into a SmartArt diagram would silently disappear from parsing, same gap
    that existed in the JS parser before extractDiagramTexts was added.
    """
    graphic_data = shape._element.find(f".//{{{NS['a']}}}graphicData")
    if graphic_data is None or "diagram" not in (graphic_data.get("uri") or ""):
        return []
    rel_ids_el = graphic_data.find(f".//{{{NS['dgm']}}}relIds")
    if rel_ids_el is None:
        return []
    dm_rel_id = rel_ids_el.get(f"{{{NS['r']}}}dm")
    if not dm_rel_id:
        return []

    slide_part = shape.part
    if dm_rel_id not in slide_part.rels:
        return []
    diagram_blob = slide_part.rels[dm_rel_id].target_part.blob
    diagram_xml = etree.fromstring(diagram_blob)

    texts = []
    for pt in diagram_xml.iter(f"{{{NS['dgm']}}}pt"):
        pt_type = pt.get("type")
        if pt_type and pt_type != "node":
            continue  # parTrans/sibTrans/doc are structural connectors, not content
        for p in pt.iter(f"{{{NS['a']}}}p"):
            text = _collapse("".join(t.text or "" for t in p.iter(f"{{{NS['a']}}}t")))
            if text:
                texts.append(text)
    return texts


def _chart_points_by_idx(container) -> dict:
    """idx -> raw <c:v> text, exactly as cached in the chart XML. Reading
    the text verbatim (rather than going through python-pptx's Series.values,
    which parses each cached value into a Python float) means a decimal like
    "4.4" is preserved exactly instead of being round-tripped through float
    parsing/re-serialization, which can surface binary floating-point
    artifacts (e.g. 4.4000000000000004) that were never in the source file.
    """
    points = {}
    if container is None:
        return points
    for pt in container.iter(f"{{{NS['c']}}}pt"):
        idx = pt.get("idx")
        v = pt.find(f"{{{NS['c']}}}v")
        if idx is not None and v is not None and v.text:
            points[idx] = v.text.strip()
    return points


def _extract_chart_texts(chart) -> list[str]:
    """Chart title + one line per series, read straight from the chart's raw
    XML (chart._chartSpace) rather than python-pptx's Chart API — see
    _chart_points_by_idx for why. Without this at all, edits to a chart's
    underlying numbers (the whole point of a chart) would be completely
    invisible to the comparison, same gap that existed in the JS parser
    before extractChartTexts was added there.
    """
    chart_space = chart._chartSpace
    texts = []

    title_el = chart_space.find(f".//{{{NS['c']}}}title")
    if title_el is not None:
        title = _collapse("".join(t.text or "" for t in title_el.iter(f"{{{NS['a']}}}t")))
        if title:
            texts.append(f"Chart title: {title}")

    for ser in chart_space.iter(f"{{{NS['c']}}}ser"):
        tx_el = ser.find(f"{{{NS['c']}}}tx")
        name_v = tx_el.find(f".//{{{NS['c']}}}v") if tx_el is not None else None
        series_name = name_v.text.strip() if name_v is not None and name_v.text else ""

        categories = _chart_points_by_idx(ser.find(f"{{{NS['c']}}}cat"))
        values = _chart_points_by_idx(ser.find(f"{{{NS['c']}}}val"))
        idxs = sorted(set(categories) | set(values), key=int)
        pairs = [f"{categories.get(idx, f'#{idx}')}: {values.get(idx, '')}" for idx in idxs]

        line = ((f"{series_name} — " if series_name else "") + ", ".join(pairs)).strip()
        if line:
            texts.append(line)
    return texts


def parse_pptx(data: bytes, filename: str) -> dict:
    prs = Presentation(io.BytesIO(data))
    sections = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        section = {"title": f"Slide {slide_num}", "page": None, "slide": slide_num, "blocks": []}
        title_shape = slide.shapes.title
        table_counter = 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                para_texts = _paragraph_texts_from_shape(shape)
                if para_texts:
                    if shape is title_shape:
                        text = " ".join(para_texts)
                        section["title"] = text
                        section["blocks"].append({"type": "heading", "text": text, "cells": None, "tableId": None, "rowIndex": None})
                    else:
                        block_type = "list" if len(para_texts) > 1 else "paragraph"
                        for text in para_texts:
                            section["blocks"].append({"type": block_type, "text": text, "cells": None, "tableId": None, "rowIndex": None})

            if shape.has_table:
                table_counter += 1
                table_id = f"slide{slide_num}_table_{table_counter}"
                for row_index, row in enumerate(shape.table.rows):
                    cells = [_collapse(cell.text) for cell in row.cells]
                    if any(cells):
                        section["blocks"].append({
                            "type": "tableRow",
                            "text": " | ".join(cells),
                            "cells": cells,
                            "tableId": table_id,
                            "rowIndex": row_index,
                        })

            # python-pptx's shape_type is documented to be None for SmartArt
            # ("This value is None when none of these [four] types apply, for
            # example when the shape contains SmartArt") — there's no reliable
            # shape_type to gate on, so just attempt extraction on every shape;
            # _extract_diagram_texts itself checks graphicData's uri and is a
            # no-op (returns []) for anything that isn't actually a diagram.
            for text in _extract_diagram_texts(shape):
                section["blocks"].append({"type": "list", "text": text, "cells": None, "tableId": None, "rowIndex": None})

            if shape.has_chart:
                for text in _extract_chart_texts(shape.chart):
                    section["blocks"].append({"type": "list", "text": text, "cells": None, "tableId": None, "rowIndex": None})

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            for paragraph in slide.notes_slide.notes_text_frame.paragraphs:
                text = _collapse(paragraph.text)
                if text:
                    section["blocks"].append({"type": "notes", "text": text, "cells": None, "tableId": None, "rowIndex": None})

        sections.append(section)

    return {"docType": "pptx", "fileName": filename, "sections": sections}
